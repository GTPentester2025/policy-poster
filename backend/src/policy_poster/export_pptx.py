"""Stage 9 — PPTX export: one slide per poster, real editable text boxes
(never flattened images). Slide dimensions per orientation:
13.33×7.5 in landscape, 7.5×13.33 in portrait (spec §4 stage 9)."""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .content import PosterContent

_LANDSCAPE = (Inches(13.33), Inches(7.5))
_PORTRAIT = (Inches(7.5), Inches(13.33))

_ACCENT = RGBColor(0x1F, 0x3A, 0x5F)
_BODY = RGBColor(0x22, 0x22, 0x22)


def _add_textbox(slide, left, top, width, height, text, size_pt,
                 bold=False, color=_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def export_pptx(content: PosterContent, orientation: str, path: str) -> str:
    if orientation not in ("landscape", "portrait"):
        raise ValueError(f"unknown orientation: {orientation!r}")
    width, height = _LANDSCAPE if orientation == "landscape" else _PORTRAIT

    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    margin = Inches(0.6)
    usable = width - 2 * margin
    y = margin

    def place(text, size, bold=False, color=_BODY, gap=0.15, lines=1):
        nonlocal y
        box_height = Inches(0.35 * lines + size / 72)
        _add_textbox(slide, margin, y, usable, box_height, text, size, bold, color)
        y = y + box_height + Inches(gap)

    place(content.eyebrow.text, 16, bold=True, color=_ACCENT)
    place(content.headline.text, 40 if orientation == "landscape" else 34,
          bold=True, lines=2)
    place(content.subhead.text, 20, lines=2)
    for point in content.body_points:
        place(f"• {point.text}", 16, gap=0.05)
    y = y + Inches(0.2)
    place(content.callout.text, 22, bold=True, color=_ACCENT)
    place(content.cta.text, 18, bold=True)

    prs.save(path)
    return path


def export_both_orientations(content: PosterContent, base_path: str) -> dict[str, str]:
    """Four-artefact minimum is completed by the JPG renders (Phase 3)."""
    out = {}
    for orientation in ("landscape", "portrait"):
        path = f"{base_path}_{orientation}.pptx"
        export_pptx(content, orientation, path)
        out[orientation] = path
    return out
