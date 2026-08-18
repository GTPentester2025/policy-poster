from pptx import Presentation
from pptx.util import Emu

from policy_poster.content import PosterContent, Slot
from policy_poster.export_pptx import export_pptx


def content():
    return PosterContent(
        poster_id="p1", angle="urgency", template_family="default",
        eyebrow=Slot("SECURITY FIRST", ["1.1"]),
        headline=Slot("Report incidents fast", ["1.1"]),
        subhead=Slot("Every incident must be reported within 24 hours.", ["1.1"]),
        body_points=[
            Slot("Records are destroyed after 90 days.", ["2.1"]),
            Slot("Escalate breaches to the CISO.", ["1.1"]),
        ],
        callout=Slot("24 hours. No exceptions.", ["1.1"]),
        cta=Slot("Report it now", ["1.1"]),
        coverage_map={"1.1": "covered"},
    )


def all_text(prs):
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def test_landscape_dimensions_and_editable_text(tmp_path):
    path = tmp_path / "poster_landscape.pptx"
    export_pptx(content(), "landscape", str(path))
    prs = Presentation(str(path))
    assert round(prs.slide_width / 914400, 2) == 13.33
    assert round(prs.slide_height / 914400, 2) == 7.5
    text = all_text(prs)
    for expected in ["SECURITY FIRST", "Report incidents fast",
                     "Records are destroyed after 90 days.", "Report it now"]:
        assert expected in text  # real, selectable text frames — not an image


def test_portrait_dimensions_same_content(tmp_path):
    lpath, ppath = tmp_path / "l.pptx", tmp_path / "p.pptx"
    export_pptx(content(), "landscape", str(lpath))
    export_pptx(content(), "portrait", str(ppath))
    prs = Presentation(str(ppath))
    assert round(prs.slide_width / 914400, 2) == 7.5
    assert round(prs.slide_height / 914400, 2) == 13.33
    # C7: byte-identical content across orientations
    assert sorted(all_text(Presentation(str(lpath))).split("\n")) == \
           sorted(all_text(prs).split("\n"))


def test_no_picture_shapes(tmp_path):
    path = tmp_path / "poster.pptx"
    export_pptx(content(), "landscape", str(path))
    prs = Presentation(str(path))
    kinds = {shape.shape_type for slide in prs.slides for shape in slide.shapes}
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    assert MSO_SHAPE_TYPE.PICTURE not in kinds
